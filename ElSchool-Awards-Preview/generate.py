# Jesse M, 18/9/24
# Last modified: ?/03/2026

# This program generates an awards slideshow starter given a CSV of
# results.

# It greatly speeds up the process of using processing results into a
# presentation.

import csv
import os
import pptx
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

# - Dev variables -
MAX_WIN_IMAGES = 4
MAX_WIN_NAMES = 2

# Presentation object
presentation = pptx.Presentation("Template.pptx")

# Setting up the CSV to be read
csv_path = [f for f in os.listdir() if f.endswith(".csv")][0]
csv_file = open(csv_path, "r")
csv_reader = csv.reader(csv_file)

total_scores = {}

# Iterate over columns in CSV (questions)
for column_index, column in enumerate(zip(*csv_reader)):
    # The question is the header of the column
    question = column[0]

    # Question scores
    question_scores = {}

    # Loop over the entries of a given column (question)
    for entry in column[1:]:
        # Format entry to votes string
        votes_str = entry.strip().title().replace(", ", ",")

        # If empty input, skip
        if votes_str == "":
            continue
        
        # Split votes into list - only take three names.
        votes = votes_str.split(",")[:3]

        for i, person in enumerate(votes):
            # The multiplier based on the tier of the vote (gold, etc)
            mult = 3 - i

            # If they don't already have score, add them to the dict
            if not person in question_scores:
                question_scores[person] = [0, 0, 0]
            
            if not person in total_scores:
                total_scores[person] = [0, 0, 0]
            
            # Increment their medal count (3 for gold, 2 for silver, etc)
            question_scores[person][i] += mult

            total_scores[person][i] += mult

    # Sorted question scores by sum, low to high
    # This is because the chart adds categories from top to bottom.
    sqs = sorted(question_scores.items(), key=lambda x: sum(x[1]))

    # Grouping people into group based on score.
    # e.g. all people with sum of 20, go in bucket 20 of the groups dict.
    groups = {}
    for person, votes in question_scores.items():
        score = sum(votes)
        if not score in groups:
            groups[score] = []
        groups[score].append(person)

    # We use these groups to yield the top 3 people, high to low.
    # Note: each rank (1st, 2nd, ...) can, and likely will have multiple people
    top3 = sorted(groups.items(), key=lambda x: x[0])[::-1]
    top3 = [x[1] for x in top3[:3]]


    # - Creating the slide -

    # Question Slide (inefficient to use two for loops, but i cba)
    slide = presentation.slides[column_index * 2]

    # Loop over shapes (things) of the slide
    for shape in slide.shapes:
        # Try and grab the object data of the shape, if not continue to next shape
        try:
            obj = shape.text_frame.paragraphs[0].runs[0]
        except:
            continue

        # The title of the slide (the question)
        if obj.text == "question":
            obj.text = question

    # Results slide
    slide = presentation.slides[column_index * 2 + 1]

    for shape in slide.shapes:
        # Try and grab the object data of the shape, if not continue to next shape
        try:
            obj = shape.text_frame.paragraphs[0].runs[0]
        except:
            continue

        # The segmented bar chart
        if obj.text == "chart":
            # Delete old shape
            shape.element.getparent().remove(shape.element)

            # Extract sorted categories (names) and their vote counts
            first_votes = [votes[0] for _, votes in sqs]
            second_votes = [votes[1] for _, votes in sqs]
            third_votes = [votes[2] for _, votes in sqs]

            # Define the chart data
            chart_data = CategoryChartData()
            chart_data.categories = [name for name, _ in sqs]
            
            # Adding series (First, Second, Third votes)
            chart_data.add_series('First Votes', first_votes)
            chart_data.add_series('Second Votes', second_votes)
            chart_data.add_series('Third Votes', third_votes)

            # Create chart
            chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.BAR_STACKED, 
                    shape.left, 
                    shape.top, 
                    shape.width, 
                    shape.height, 
                    chart_data
            ).chart
            
            # Define the colors for each series
            colors = [
                RGBColor(255, 215, 0),    # Gold (First votes)
                RGBColor(192, 192, 192),  # Silver (Second votes)
                RGBColor(205, 127, 50)    # Bronze (Third votes)
            ]
            # Apply colors to the series (each segment of the bar)
            for i, series in enumerate(chart.series):
                for point in series.points:
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = colors[i]

            # Set the colour of the chart to black
            chart.category_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black font colour
            chart.value_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black font colour

        # A winner image
        elif obj.text.startswith("img"):
            # Delete old shape
            shape.element.getparent().remove(shape.element)

            # Extract the winners for this rank, if there are any
            # Limited to as many as we can fit on the actual slideshow
            try:
                winners = top3[int(obj.text[-1]) - 1][:MAX_WIN_IMAGES]
            except IndexError:
                winners = ["N/A"]

            # Fit all winners into the single image space
            for i, winner in enumerate(winners):
                img_width = shape.width / len(winners)

                winner_img_path = f"Images/{winner}.png"

                # If image doesn't exist, resort to deault
                img_path = winner_img_path if os.path.isfile(winner_img_path) else "Images/NoImage.png"

                # Add image
                img_element = slide.shapes.add_picture(
                    img_path, shape.left + img_width * i, shape.top,
                    img_width, shape.height
                )
                # Send image to the back of slide
                img_element._element.getparent().insert(0, img_element._element)

        # The winner text (above each rank)
        elif obj.text.startswith("win"):
            # Extract the winners for this rank - all winners, if applicable
            try:
                winners = top3[int(obj.text[-1]) - 1]
            except IndexError:
                winners = ["N/A"]

            # If there are two many winners, truncate and add '...'
            if len(winners) > MAX_WIN_NAMES:
                winners = winners[:MAX_WIN_NAMES] + ["..."]

            # Modify the text
            obj.text = ", ".join(winners)

        # The title of the slide (the question)
        elif obj.text == "question":
            obj.text = question


# Sorted total scores, high to low
sts = sorted(total_scores.items(), key=lambda x: sum(x[1]))[::-1]

# Printing the total scores for copy + paste
print("Total Scores:\n")
for person, votes in sts:
    print(f"{person} - {sum(votes)}")


# Close CSV once finished using the CSV reader object.
csv_file.close()

# Save presentation
presentation.save("New-Slideshow.pptx")

print()
input("Press enter to close.")
